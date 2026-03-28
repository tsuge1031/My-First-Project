from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import subprocess

# ── カラー定義 ──────────────────────────────
DARK_BLUE   = (30,  58,  95)
MID_BLUE    = (42,  82,  152)
ACCENT      = (0,   119, 204)
GREEN       = (10,  124, 66)
RED         = (192, 57,  43)
ORANGE      = (211, 84,  0)
WHITE       = (255, 255, 255)
LIGHT_GRAY  = (243, 244, 246)
LIGHT_BLUE  = (232, 244, 253)
LIGHT_GREEN = (230, 246, 238)
LIGHT_RED   = (253, 240, 238)
TEXT_DARK   = (26,  32,  44)
TEXT_SUB    = (74,  85,  104)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── ユーティリティ ────────────────────────────
def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_rect(slide, left, top, width, height, fill_color, line=False, line_color=None):
    s = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(*fill_color)
    if line and line_color:
        s.line.color.rgb = RGBColor(*line_color)
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def add_text(slide, left, top, width, height, text, size, bold=False,
             color=WHITE, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txb.text_frame
    tf.word_wrap = wrap
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name      = "Meiryo"
        run.font.size      = Pt(size)
        run.font.bold      = bold
        run.font.italic    = italic
        run.font.color.rgb = RGBColor(*color)

def add_card(slide, left, top, width, height, title, body,
             bg=LIGHT_GRAY, title_color=TEXT_DARK, body_color=TEXT_SUB,
             title_size=13, body_size=11):
    add_rect(slide, left, top, width, height, bg)
    if title:
        add_text(slide, left+0.12, top+0.1, width-0.24, 0.35,
                 title, title_size, bold=True, color=title_color)
    if body:
        add_text(slide, left+0.12, top+(0.42 if title else 0.12),
                 width-0.24, height-(0.52 if title else 0.24),
                 body, body_size, color=body_color)

def add_section_title(slide, text, sub=""):
    add_rect(slide, 0, 0, 13.33, 1.1, DARK_BLUE)
    add_text(slide, 0.5, 0.15, 12, 0.6, text, 24, bold=True, color=WHITE)
    if sub:
        add_text(slide, 0.5, 0.72, 12, 0.35, sub, 12, color=(180,200,230))

# ════════════════════════════════════════════
# 1. 表紙
# ════════════════════════════════════════════
slide = blank_slide(prs)
add_rect(slide, 0, 0, 13.33, 7.5, DARK_BLUE)
add_rect(slide, 0, 5.8, 13.33, 1.7, MID_BLUE)
add_rect(slide, 0.5, 1.6, 0.08, 2.8, ACCENT)  # 左アクセントライン

add_text(slide, 0.8, 0.6,  12, 0.5,  "役員向け・社外秘", 12, color=(160,185,220))
add_text(slide, 0.8, 1.3,  12, 1.1,  "Claude Code 社内導入提案", 38, bold=True, color=WHITE)
add_text(slide, 0.8, 2.5,  12, 0.5,  "〜 AIエージェントハブとして、全業務を変える 〜", 16, color=(120,170,220))
add_text(slide, 0.8, 3.2,  12, 1.2,
         "コードを書くツールではない。あらゆるAIサービス・社内システムと連携し、\n"
         "開発・営業・管理・提案まで横断的に業務を自動化する次世代AI基盤の導入提案。",
         13, color=(180,200,230))
add_text(slide, 0.8, 5.95, 6,  0.45, "2026年3月28日", 12, color=(200,215,235))
add_text(slide, 6.5, 5.95, 6,  0.45, "対象：全部門・役員　｜　機密区分：社内限り",
         12, color=(200,215,235), align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════
# 2. Claude Code とは何か（ハブ図解）
# ════════════════════════════════════════════
slide = blank_slide(prs)
add_section_title(slide, "01｜Claude Code とは何か", "「コードを書くAI」ではなく「AIエージェントハブ」")

# 中央ハブ
add_rect(slide,  5.67, 2.8, 2.0, 2.0, MID_BLUE)
add_text(slide,  5.67, 2.9, 2.0, 0.5, "🤖", 26, align=PP_ALIGN.CENTER)
add_text(slide,  5.67, 3.4, 2.0, 0.5, "Claude Code", 13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide,  5.67, 3.85, 2.0, 0.4, "AIエージェントハブ", 10, color=(180,210,240), align=PP_ALIGN.CENTER)

# 左側ノード
left_nodes = [
    ("📧", "メール / カレンダー", "Gmail / Outlook", LIGHT_BLUE, ACCENT),
    ("📋", "プロジェクト管理",    "Jira / Notion",   LIGHT_GREEN, GREEN),
    ("💬", "コミュニケーション", "Slack / Teams",   (240,235,255), (109,40,217)),
    ("🗄️", "社内DB・基幹システム", "PostgreSQL / SAP", (255,245,235), ORANGE),
]
for i, (icon, title, sub, bg, tc) in enumerate(left_nodes):
    y = 1.35 + i * 1.35
    add_card(slide, 0.3, y, 4.2, 1.1, f"{icon} {title}", sub, bg, tc, TEXT_SUB, 12, 10)
    add_text(slide, 4.5, y+0.35, 1.2, 0.4, "⇄", 18, color=ACCENT, align=PP_ALIGN.CENTER)

# 右側ノード
right_nodes = [
    ("🔎", "Web検索・情報収集", "Google / Bing",    LIGHT_BLUE, ACCENT),
    ("💻", "開発ツール",        "GitHub / CI/CD",   LIGHT_GREEN, GREEN),
    ("💳", "決済・会計",        "Stripe / 弥生",    (255,245,235), ORANGE),
    ("🌐", "ブラウザ自動操作",  "Playwright MCP",   (240,235,255), (109,40,217)),
]
for i, (icon, title, sub, bg, tc) in enumerate(right_nodes):
    y = 1.35 + i * 1.35
    add_card(slide, 8.83, y, 4.2, 1.1, f"{icon} {title}", sub, bg, tc, TEXT_SUB, 12, 10)
    add_text(slide, 7.63, y+0.35, 1.2, 0.4, "⇄", 18, color=ACCENT, align=PP_ALIGN.CENTER)

# MCP ラベル
add_rect(slide, 3.5, 6.75, 6.33, 0.5, MID_BLUE)
add_text(slide, 3.5, 6.77, 6.33, 0.45,
         "MCP（Model Context Protocol）= AIと外部システムをつなぐ業界標準プロトコル",
         11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# 3. 業務活用領域
# ════════════════════════════════════════════
slide = blank_slide(prs)
add_section_title(slide, "02｜活用できる業務領域", "開発部門だけでなく、全部門の業務を変える")

cards = [
    ("📝 提案・営業支援",    "顧客情報・過去案件・市場データを横断収集し、提案書・見積書を自動ドラフト", "営業部門", LIGHT_BLUE,  ACCENT),
    ("💻 開発業務の自動化",  "コード生成・レビュー・テスト・ドキュメント作成を自動化。指示一つで実行",  "開発部門", LIGHT_GREEN, GREEN),
    ("📊 報告書・議事録",    "会議メモをSlack/Notionから取得し、議事録・月次報告書を自動生成",         "全部門",   (240,235,255), (109,40,217)),
    ("🔄 業務フロー自動化",  "申請→承認→通知→記録の一連フローをシステム横断で自動実行・無人化",         "管理部門", (255,245,235), ORANGE),
    ("🔍 情報収集・競合分析", "Web・業界サイトを自律的に収集・整理し、リサーチレポートを自動生成",      "企画・マーケ", LIGHT_RED, RED),
    ("🤝 顧客向けAI導入支援", "自社導入ノウハウを活かしAIエージェント導入支援を新サービスとして提供",  "新規事業",  LIGHT_GRAY, DARK_BLUE),
]
cols, col_w, col_h = 3, 4.1, 2.1
for i, (title, body, tag, bg, tc) in enumerate(cards):
    col, row = i % cols, i // cols
    x = 0.3 + col * (col_w + 0.2)
    y = 1.25 + row * (col_h + 0.15)
    add_card(slide, x, y, col_w, col_h, title, body + f"\n\n  [{tag}]", bg, tc, TEXT_SUB, 13, 11)

# ════════════════════════════════════════════
# 4. 国内導入実績
# ════════════════════════════════════════════
slide = blank_slide(prs)
add_section_title(slide, "03｜国内企業の導入実績（2025〜2026年）", "コード以外の業務でも大きな成果が出ている")

cases = [
    ("Findy（採用プラットフォーム）",
     "1時間 → 5分",
     "求人票作成時間の短縮",
     "Notion MCPで企業情報を自動取得し、求人票ドラフトを自動生成。\nエンジニア以外のビジネス職が独立して利用できるカスタムコマンドとして実装。品質のばらつきも解消。",
     "Notion MCP"),
    ("LINEヤフー",
     "週6時間削減",
     "レビュー準備の工数削減",
     "GitHub MCP・Jira MCP・Confluence MCPを連携。\nPRレビューに必要な設計書・企画書・関連タスクの情報収集を全自動化。手作業ゼロでレビュー準備が完結。",
     "GitHub / Jira / Confluence MCP"),
    ("Room8（中小企業DX事例）",
     "150万円 → 700円",
     "業務自動化の実現コスト",
     "申込→決済→サブスク設定→入金確認までの業務フローをClaude Code × Stripe APIで完全無人化。\n月額500円のVPS上で運用し、従来150万円の開発コストを大幅削減。",
     "Stripe API MCP"),
]
for i, (company, result, unit, desc, tool) in enumerate(cases):
    x = 0.3 + i * 4.35
    add_rect(slide, x, 1.25, 4.1, 5.8, LIGHT_GRAY, True, (210,220,230))
    add_text(slide, x+0.15, 1.35, 3.8, 0.4, company, 10, bold=True, color=TEXT_SUB)
    add_text(slide, x+0.15, 1.8,  3.8, 0.9, result, 28, bold=True, color=MID_BLUE)
    add_text(slide, x+0.15, 2.65, 3.8, 0.4, unit,   12, color=TEXT_SUB)
    add_rect(slide, x+0.15, 3.05, 3.8, 0.02, (210,220,230))
    add_text(slide, x+0.15, 3.15, 3.8, 2.4, desc, 11, color=TEXT_SUB)
    add_rect(slide, x+0.15, 6.6, 3.8, 0.35, WHITE, True, (200,215,235))
    add_text(slide, x+0.2, 6.62, 3.7, 0.32, f"🔧 {tool}", 10, bold=True, color=MID_BLUE)

# ════════════════════════════════════════════
# 5. メリット・デメリット
# ════════════════════════════════════════════
slide = blank_slide(prs)
add_section_title(slide, "04｜メリット・デメリット", "")

merits = [
    ("全業務部門の生産性向上",    "開発・営業・管理・企画を問わず対象。定型業務・情報収集・文書作成をAIが代替"),
    ("既存システムとの即時連携",  "MCPにより社内DB・Slack・Jira等と接続。新規開発不要で既存IT資産をAIから活用"),
    ("顧客への新サービス展開",    "導入ノウハウを武器にAIエージェント導入支援・MCP構築を顧客向けサービスとして提供"),
    ("業界標準技術の先行習得",    "MCPはOpenAI・Googleも採用する標準規格。先行習熟で技術的優位性を確立"),
    ("競争力・差別化の実現",      "AI活用実績を提案力に変換。工数削減による価格競争力と技術力ブランドを両立"),
]
demerits = [
    ("顧客情報の外部送信リスク",  "社内DBへのMCP接続が強力なため設定ミスによる情報流出リスクあり。権限制御が必須"),
    ("エージェントの暴走リスク",  "自律動作による意図しない操作（ファイル削除・メール送信等）に注意。承認フロー設計が重要"),
    ("MCP設計の専門知識が必要",   "MCPサーバー構築・管理にエンジニアリング知識が必要。担当者の育成が必要"),
    ("ランニングコストの変動",    "複数システム呼び出しでAPI使用量増加。利用上限設定と定期的なコスト監視が必要"),
    ("AI出力への過信リスク",      "自律実行の便利さから人間のチェックが薄れる危険性。重要業務は最終確認ルールを設計"),
]

add_rect(slide, 0.3, 1.2, 6.2, 0.55, GREEN)
add_text(slide, 0.4, 1.22, 6.0, 0.5, "✅  メリット（導入効果）", 14, bold=True, color=WHITE)
for i, (title, body) in enumerate(merits):
    add_card(slide, 0.3, 1.85 + i*1.05, 6.2, 0.95, f"▶ {title}", body, LIGHT_GREEN, GREEN, TEXT_SUB, 12, 10)

add_rect(slide, 6.83, 1.2, 6.2, 0.55, RED)
add_text(slide, 6.93, 1.22, 6.0, 0.5, "⚠️  デメリット・留意点", 14, bold=True, color=WHITE)
for i, (title, body) in enumerate(demerits):
    add_card(slide, 6.83, 1.85 + i*1.05, 6.2, 0.95, f"▶ {title}", body, LIGHT_RED, RED, TEXT_SUB, 12, 10)

# ════════════════════════════════════════════
# 6. セキュリティ対策
# ════════════════════════════════════════════
slide = blank_slide(prs)
add_section_title(slide, "05｜セキュリティ対策", "独立系SIerとして最重要——技術的に担保できる")

sec_cards = [
    ("🔒 学習利用なし",          "Team/Enterprise/APIプランでは送信データがモデル学習に一切使用されないことを契約で保証"),
    ("🗄️ ゼロデータ保持（ZDR）",  "Enterpriseプランで申請可能。推論後にサーバーから即時削除。官公庁・金融案件にも対応"),
    ("☁️ 自社クラウド内完結",    "AWS Bedrock / Google Vertex AI経由でデータがAnthropicサーバーに送られず自社環境内で完結"),
    ("⚙️ MCPアクセス権限制御",   "MCP接続先ごとに読み取り専用・書き込み禁止を細かく設定。最小権限の原則を実装"),
    ("📋 監査ログ完備",           "AWS CloudTrailによる自動記録。AIの全アクションを追跡・証跡化。コンプライアンス対応"),
    ("🛡️ 国際認証取得済み",      "SOC 2 Type II・ISO 27001・ISO/IEC 42001（AI管理）・HIPAA BAA対応。第三者審査済み"),
]
cols3, w3, h3 = 3, 4.1, 1.9
for i, (title, body) in enumerate(sec_cards):
    col, row = i % cols3, i // cols3
    x = 0.3 + col * (w3 + 0.2)
    y = 1.25 + row * (h3 + 0.2)
    add_card(slide, x, y, w3, h3, title, body, LIGHT_GRAY, MID_BLUE, TEXT_SUB, 12, 11)

# 推奨方式バナー
add_rect(slide, 0.3, 5.45, 12.73, 0.55, MID_BLUE)
add_text(slide, 0.5, 5.47, 12.5, 0.5,
         "推奨：AWS Bedrock経由　→　データが自社AWS内に留まり CloudTrail で監査ログを完備",
         13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(slide, 0.3, 6.1, 5.9, 0.8, LIGHT_BLUE, True, (150,195,230))
add_text(slide, 0.4, 6.12, 5.7, 0.75,
         "Anthropic API直接 → PoC検証用のみ推奨", 11, color=MID_BLUE)
add_rect(slide, 6.5, 6.1, 5.9, 0.8, (240,255,245), True, (150,210,175))
add_text(slide, 6.6, 6.12, 5.7, 0.75,
         "Enterprise + ZDR → 官公庁・金融など最高水準が求められる案件向け", 11, color=GREEN)

# ════════════════════════════════════════════
# 7. 導入ロードマップ
# ════════════════════════════════════════════
slide = blank_slide(prs)
add_section_title(slide, "06｜導入ロードマップ", "段階的に進めてリスクを最小化")

phases = [
    ("Phase 1：PoC検証", "期間：1〜2ヶ月", [
        "開発5〜10名・非開発2〜3名で試験導入",
        "AWS Bedrock経由のセキュア環境構築",
        "MCP：Slack / Notion / GitHubから開始",
        "利用ガイドライン・権限設計",
        "業務別の効果を定量計測",
    ], "文書作成・情報収集工数 30%削減", MID_BLUE),
    ("Phase 2：部門展開", "期間：2〜4ヶ月", [
        "全部門に展開・業務別MCP整備",
        "社内DB・基幹システムとのMCP連携",
        "部門別カスタムコマンド整備",
        "マルチエージェント活用の試験運用",
        "社員教育・活用推進担当の設置",
    ], "定型業務 50%自動化・残業削減", GREEN),
    ("Phase 3：外部展開・収益化", "期間：5ヶ月以降", [
        "顧客向けAIエージェント導入支援",
        "MCP構築サービスの提供開始",
        "AI活用実績を営業ツールとして活用",
        "提案書への「AI活用実績」訴求追加",
        "継続的なコスト最適化",
    ], "新規AIサービス売上の創出", ORANGE),
]
pw = (13.33 - 0.9) / 3
for i, (title, period, items, kpi, color) in enumerate(phases):
    x = 0.3 + i * (pw + 0.15)
    add_rect(slide, x, 1.25, pw, 0.65, color)
    add_text(slide, x+0.1, 1.27, pw-0.2, 0.6, title, 14, bold=True, color=WHITE)
    add_rect(slide, x, 1.95, pw, 4.2, LIGHT_GRAY)
    add_text(slide, x+0.12, 2.02, pw-0.24, 0.35, period, 11, color=TEXT_SUB, italic=True)
    body = "\n".join(f"• {item}" for item in items)
    add_text(slide, x+0.12, 2.42, pw-0.24, 3.6, body, 12, color=TEXT_DARK)
    add_rect(slide, x, 6.2, pw, 0.75, (220, 230, 245))
    add_text(slide, x+0.12, 6.22, pw-0.24, 0.3, "KPI目標", 10, bold=True, color=MID_BLUE)
    add_text(slide, x+0.12, 6.5,  pw-0.24, 0.4, kpi, 11, color=TEXT_DARK)

# ════════════════════════════════════════════
# 8. まとめ・経営への提言
# ════════════════════════════════════════════
slide = blank_slide(prs)
add_section_title(slide, "07｜まとめ・経営への提言", "")

# 数値サマリー（横4列）
summaries = [
    ("全部門", "活用対象", "開発・営業・管理・企画"),
    ("MCP", "業界標準の接続技術", "OpenAI・Googleも採用"),
    ("約21倍", "月次ROI（試算）", "20名・生産性20%向上ベース"),
    ("新収益", "顧客への展開可能", "AI導入支援サービス化"),
]
sw = (13.33 - 0.8) / 4
for i, (num, unit, label) in enumerate(summaries):
    x = 0.4 + i * (sw + 0.0)
    add_rect(slide, x, 1.2, sw-0.15, 1.4, LIGHT_GRAY, True, (210,220,230))
    add_text(slide, x, 1.25, sw-0.15, 0.7, num,   22, bold=True, color=MID_BLUE,  align=PP_ALIGN.CENTER)
    add_text(slide, x, 1.88, sw-0.15, 0.4, unit,  11, bold=False, color=TEXT_SUB, align=PP_ALIGN.CENTER)
    add_text(slide, x, 2.28, sw-0.15, 0.28, label, 9, color=TEXT_SUB, align=PP_ALIGN.CENTER)

# 提言ボックス
add_rect(slide, 0.3, 2.8, 12.73, 4.3, DARK_BLUE)
add_text(slide, 0.55, 2.92, 12, 0.5, "🎯  経営判断のポイント", 16, bold=True, color=WHITE)

points = [
    ("1", "「コードを書くAI」ではなく「業務基盤としてのAI」",
     "Claude Codeは開発部門だけのツールではない。MCPによって全社の業務システムと接続し、全部門の生産性を底上げする基盤投資として捉えるべき。"),
    ("2", "シャドーAIを放置することの方がリスクが大きい",
     "社員が個人でAIツールを無断利用している可能性が高い。公式導入によりガバナンスの下に置く方が、顧客情報漏洩リスクを確実に低減できる。"),
    ("3", "MCP習熟は今後の受注競争力に直結する",
     "MCPはOpenAI・Googleも採用する業界標準。今から自社実装することで、顧客のAIシステム統合案件を受注する技術的優位性を早期確立できる。"),
    ("4", "推奨アクション：まずPoC承認を",
     "AWS Bedrock経由のセキュアな環境で1〜2ヶ月のPoC実施を承認いただきたい。定量効果確認後に全社展開・顧客サービス化を判断する。"),
]
for i, (num, title, body) in enumerate(points):
    y = 3.55 + i * 0.88
    add_rect(slide, 0.45, y, 0.38, 0.38, MID_BLUE)
    add_text(slide, 0.45, y, 0.38, 0.38, num, 13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, 0.9, y,      11.9, 0.32, title, 13, bold=True,  color=WHITE)
    add_text(slide, 0.9, y+0.33, 11.9, 0.52, body,  10, bold=False, color=(180,200,230))

# ════════════════════════════════════════════
# 保存・起動
# ════════════════════════════════════════════
output = "/Users/tsuge/Desktop/My-First-Project/20260328_ClaudeCode社内導入提案_役員向け.pptx"
prs.save(output)
print(f"✅ 保存完了: {output}")
subprocess.run(["open", output])
